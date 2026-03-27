"""Build NVC pitch deck for Parachute — 5 minute pitch.
9 slides, brand palette, consistent layout."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Brand colors ──
BG       = RGBColor(0xFA, 0xF8, 0xF4)
BG_SOFT  = RGBColor(0xF3, 0xF0, 0xEA)
FG       = RGBColor(0x2C, 0x2A, 0x26)
FG_MUTED = RGBColor(0x6B, 0x68, 0x60)
FG_DIM   = RGBColor(0x9A, 0x96, 0x90)
ACCENT   = RGBColor(0x4A, 0x7C, 0x59)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BORDER   = RGBColor(0xE4, 0xE0, 0xD8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

# ── Font choices (safe for presentation machines) ──
FONT_SERIF = "Georgia"
FONT_SANS = "Calibri"
FONT_MONO = "Courier New"


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_brand_bar(slide):
    """Thin accent stripe at top of slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_footer(slide, text="parachute.computer"):
    """Subtle footer with URL."""
    txBox = slide.shapes.add_textbox(
        Inches(1.2), Inches(7.0), Inches(4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = FG_DIM
    p.font.name = FONT_SANS


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=FG, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_SANS, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    return tf


def section_label(slide, text, top=Inches(1.2)):
    add_text_box(slide, Inches(1.2), top, Inches(4), Inches(0.4),
                 text.upper(), font_size=11, color=ACCENT, bold=True,
                 font_name=FONT_SANS)


def slide_headline(slide, text, top=Inches(1.7), size=40):
    add_text_box(slide, Inches(1.2), top, Inches(10), Inches(1.5),
                 text, font_size=size, color=FG, bold=False,
                 font_name=FONT_SERIF, line_spacing=1.05)


def new_slide():
    """Create a new slide with standard background, brand bar, and footer."""
    s = prs.slides.add_slide(blank_layout)
    set_slide_bg(s)
    add_brand_bar(s)
    add_footer(s)
    return s


def remove_table_borders(table):
    """Make table borders thin and subtle."""
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.find(qn(border_name))
                if ln is None:
                    ln = tcPr.makeelement(qn(border_name), {})
                    tcPr.append(ln)
                ln.set('w', '6350')  # 0.5pt
                solidFill = ln.find(qn('a:solidFill'))
                if solidFill is None:
                    solidFill = ln.makeelement(qn('a:solidFill'), {})
                    ln.append(solidFill)
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is None:
                    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {})
                    solidFill.append(srgbClr)
                srgbClr.set('val', 'E4E0D8')


# ═══════════════════════════════════════════════
# SLIDE 1 — The Landscape
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Landscape")
slide_headline(s, "Everyone is building\npersonal AI computers.")

add_text_box(s, Inches(1.2), Inches(3.5), Inches(10), Inches(1.0),
             "Over 100 million people already pay $20+/month for AI.\n"
             "OpenClaw. Claude Cowork. ZoComputer. Perplexity Computer. Manus.",
             font_size=18, color=FG_MUTED, line_spacing=1.5)

# Market stat
add_text_box(s, Inches(1.2), Inches(5.0), Inches(3), Inches(1.2),
             "$50B+", font_size=64, color=ACCENT, font_name=FONT_SERIF)
add_text_box(s, Inches(4.5), Inches(5.25), Inches(5), Inches(0.8),
             "agentic AI market by 2030 · 44% CAGR",
             font_size=16, color=FG_MUTED, line_spacing=1.4)

# The turn
add_text_box(s, Inches(1.2), Inches(6.3), Inches(10), Inches(0.6),
             "But they\u2019re building walled gardens \u2014 and only for the 5% who are already power users.",
             font_size=22, color=FG, bold=False, font_name=FONT_SERIF)


# ═══════════════════════════════════════════════
# SLIDE 2 — The 95% Gap
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Opportunity")
slide_headline(s, "The other 95% just want\nsomething that works.")

add_text_box(s, Inches(1.2), Inches(3.5), Inches(10), Inches(1.5),
             "The artist who wants to organize creative ideas.\n"
             "The small business owner tracking their days.\n"
             "The parent who wants a better way to remember and reflect.",
             font_size=22, color=FG_MUTED, line_spacing=1.6, font_name=FONT_SERIF)

add_text_box(s, Inches(1.2), Inches(5.7), Inches(10), Inches(1.0),
             "They need a system that learns them \u2014\n"
             "not one that demands they learn it.",
             font_size=22, color=FG, line_spacing=1.5, font_name=FONT_SERIF)


# ═══════════════════════════════════════════════
# SLIDE 3 — Parachute Computer
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Our Answer")
slide_headline(s, "Parachute Computer.\nThe personal AI you can trust.")

# Slim bullet points with accent dot
bullets = [
    "Open source (AGPL-3.0) \u00b7 local-first \u00b7 self-hostable",
    "Your data stays on your device \u2014 portable and exportable",
    "Knowledge graph connects your journals, conversations, and thinking",
    "Public Benefit Corporation \u2014 legally mandated to serve users",
]

y = Inches(3.5)
for b in bullets:
    # Accent dot
    add_text_box(s, Inches(1.3), y + Inches(0.02), Inches(0.3), Inches(0.35),
                 "\u2022", font_size=16, color=ACCENT)
    add_text_box(s, Inches(1.7), y, Inches(9), Inches(0.4),
                 b, font_size=18, color=FG_MUTED, line_spacing=1.3)
    y += Inches(0.55)

add_text_box(s, Inches(1.2), Inches(5.8), Inches(10), Inches(0.8),
             "Software gets cloned in a day.\nTrust and context cannot.",
             font_size=22, color=FG, font_name=FONT_SERIF, line_spacing=1.3)

add_text_box(s, Inches(1.2), Inches(6.6), Inches(10), Inches(0.5),
             "But most people don\u2019t want to self-host a server.",
             font_size=18, color=FG_MUTED, font_name=FONT_SERIF)


# ═══════════════════════════════════════════════
# SLIDE 4 — Parachute Daily
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Bridge")
slide_headline(s, "Parachute Daily.\nJust talk.")

add_text_box(s, Inches(1.2), Inches(3.5), Inches(5.5), Inches(2.8),
             "A voice-first journal.\n\n"
             "Speak into a wearable pendant or your phone \u2014\n"
             "on a walk, in the car, wherever thinking happens.\n\n"
             "AI weaves through gently:\n"
             "  \u00b7  Daily reflections on your entries\n"
             "  \u00b7  Pattern recognition over time\n"
             "  \u00b7  Weekly synthesis of your thinking",
             font_size=18, color=FG_MUTED, line_spacing=1.5)

add_text_box(s, Inches(1.2), Inches(6.5), Inches(5.5), Inches(0.5),
             "No learning curve. No technical sophistication required.",
             font_size=18, color=FG, line_spacing=1.2)

# Pendant callout
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(7.8), Inches(3.5), Inches(4.5), Inches(3.0))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_SOFT
shape.line.color.rgb = BORDER
shape.line.width = Pt(1)

add_text_box(s, Inches(8.2), Inches(3.8), Inches(3.7), Inches(0.4),
             "THE PENDANT", font_size=11, color=ACCENT, bold=True)
add_text_box(s, Inches(8.2), Inches(4.3), Inches(3.7), Inches(2.0),
             "Wearable voice capture device.\n\n"
             "Go for a walk. Talk.\n"
             "Your thoughts become structured\n"
             "knowledge by the time you\u2019re home.\n\n"
             "Working prototype on stage today.",
             font_size=16, color=FG_MUTED, line_spacing=1.45)


# ═══════════════════════════════════════════════
# SLIDE 5 — Context Compounds
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Insight")
slide_headline(s, "Context compounds.")

add_text_box(s, Inches(1.2), Inches(3.3), Inches(10), Inches(0.8),
             "Every day you use Parachute, your AI gets meaningfully better.\n"
             "No competitor can shortcut months of accumulated personal context.",
             font_size=20, color=FG_MUTED, line_spacing=1.5)

# Flywheel steps with left border accent
steps = [
    ("01", "Start journaling in Daily"),
    ("02", "Build months of personal context"),
    ("03", "Upgrade to Parachute Computer"),
    ("04", "Brain ingests your history instantly"),
    ("\u2192",  "System already knows you"),
]

y = Inches(4.6)
for num, text in steps:
    is_final = num == "\u2192"
    col = ACCENT if is_final else FG
    num_col = ACCENT if is_final else FG_DIM

    add_text_box(s, Inches(1.5), y, Inches(0.6), Inches(0.45),
                 num, font_size=14, color=num_col, font_name=FONT_MONO)
    add_text_box(s, Inches(2.2), y, Inches(6), Inches(0.45),
                 text, font_size=20, color=col, font_name=FONT_SERIF,
                 bold=is_final)
    y += Inches(0.5)


# ═══════════════════════════════════════════════
# SLIDE 6 — Business Model / Pricing (5 tiers)
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Business Model")
slide_headline(s, "Start free.\nGrow with every user.")

tiers = [
    ("Free",    "Offline journal + on-device transcription. Zero hosting cost.",         FG_DIM),
    ("$2/mo",   "Cloud sync across devices. Your notes everywhere.",                     FG_MUTED),
    ("$5/mo",   "Cloud transcription + cleanup \u2014 server-side, better accuracy",     FG_MUTED),
    ("$10/mo",  "AI reflections, synthesis, and pattern surfacing",                      ACCENT),
    ("$40/mo",  "Hosted Parachute Computer \u2014 full agentic platform with bundled AI", FG),
]

y = Inches(3.4)
tier_spacing = Inches(0.6)  # Tighter for 5 tiers
for price, desc, col in tiers:
    is_highlight = price == "$10/mo"
    if is_highlight:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.0), y - Inches(0.1),
                                   Inches(10.5), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xEC)
        shape.line.fill.background()

    add_text_box(s, Inches(1.2), y, Inches(1.8), Inches(0.4),
                 price, font_size=22, color=col, font_name=FONT_SERIF)
    add_text_box(s, Inches(3.2), y + Inches(0.03), Inches(8), Inches(0.4),
                 desc, font_size=15, color=FG_MUTED)
    y += tier_spacing

add_text_box(s, Inches(1.2), Inches(6.5), Inches(10), Inches(0.6),
             "Free tier has zero hosting cost. AI at $10/mo uses cost-efficient models.\n"
             "Margins are real and improve as model costs drop.",
             font_size=14, color=FG_DIM, line_spacing=1.5)


# ═══════════════════════════════════════════════
# SLIDE 7 — Financial Projections
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Financial Projections")
slide_headline(s, "Profitable by year three.")

headers = ["", "2026", "2027", "2028"]
rows = [
    ["Free + sync users", "5,000", "50,000", "250,000"],
    ["Paid subscribers", "500", "5,000", "25,000"],
    ["Avg rev / subscriber", "~$7/mo", "~$9/mo", "~$12/mo"],
    ["ARR", "~$42K", "~$540K", "~$3.6M"],
    ["Team costs", "~$150K", "~$400K", "~$800K"],
    ["Infra + COGS", "~$15K", "~$130K", "~$500K"],
    ["Total opex", "~$165K", "~$530K", "~$1.3M"],
]

col_widths = [Inches(2.8), Inches(1.8), Inches(1.8), Inches(1.8)]
table_left = Inches(1.5)
table_top = Inches(3.3)

tbl_shape = s.shapes.add_table(len(rows) + 1, 4, table_left, table_top,
                                sum(col_widths), Inches(3.0))
tbl = tbl_shape.table

for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

# Style header
for i, h in enumerate(headers):
    cell = tbl.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.color.rgb = FG_DIM
        p.font.bold = True
        p.font.name = FONT_SANS
        p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
    cell.fill.solid()
    cell.fill.fore_color.rgb = BG_SOFT

# Style data rows
for r, row_data in enumerate(rows):
    for c, val in enumerate(row_data):
        cell = tbl.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.name = FONT_SANS
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r == 3 and c > 0:  # ARR row highlight
                p.font.color.rgb = ACCENT
                p.font.bold = True
            elif c == 0:
                p.font.color.rgb = FG
                p.font.bold = True
                p.font.size = Pt(11)
            else:
                p.font.color.rgb = FG_MUTED
        # Alternating row colors
        cell.fill.solid()
        if r == 3:  # ARR row gets light green
            cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xF2)
        elif r % 2 == 0:
            cell.fill.fore_color.rgb = WHITE
        else:
            cell.fill.fore_color.rgb = BG

# Subtle borders
remove_table_borders(tbl)

# Bottom note
add_text_box(s, Inches(1.5), Inches(6.5), Inches(9), Inches(0.4),
             "Year two approaches breakeven. Year three clearly profitable as tier mix shifts upward.",
             font_size=13, color=FG_DIM)


# ═══════════════════════════════════════════════
# SLIDE 8 — Self-funded + Team
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "Team & Traction")
slide_headline(s, "Self-funded.\nBuilt from scratch.")

# What exists — left side with accent bullets
items = [
    "Working app with local voice transcription + graph-native memory",
    "Multi-platform: phone, Telegram, Discord, Matrix",
    "Functional pendant prototype (on stage today)",
    "Daily beta launching this month \u00b7 PBC incorporated in Colorado",
]

y = Inches(3.6)
for item in items:
    add_text_box(s, Inches(1.3), y + Inches(0.02), Inches(0.3), Inches(0.3),
                 "\u2022", font_size=14, color=ACCENT)
    add_text_box(s, Inches(1.7), y, Inches(5), Inches(0.4),
                 item, font_size=15, color=FG_MUTED, line_spacing=1.3)
    y += Inches(0.5)

# Team on the right in a card
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(7.5), Inches(3.3), Inches(5), Inches(3.2))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_SOFT
shape.line.color.rgb = BORDER
shape.line.width = Pt(1)

add_text_box(s, Inches(7.9), Inches(3.5), Inches(4), Inches(0.35),
             "THE TEAM", font_size=11, color=ACCENT, bold=True)

team = [
    ("Aaron Gabriel Neyer", "Founder \u00b7 Product & architecture"),
    ("Jon Bo", "Daily co-lead \u00b7 Founding engineer"),
    ("Lucian Hymer", "Computer co-lead \u00b7 Founding engineer"),
    ("Marvin Melzer", "Hardware \u00b7 Pendant prototype"),
    ("Neil Yarnal", "Brand & design"),
]

y = Inches(4.0)
for name, role in team:
    add_text_box(s, Inches(7.9), y, Inches(4.5), Inches(0.28),
                 name, font_size=13, color=FG, bold=True)
    add_text_box(s, Inches(7.9), y + Inches(0.25), Inches(4.5), Inches(0.28),
                 role, font_size=11, color=FG_MUTED)
    y += Inches(0.5)

# Founder bio at bottom
add_text_box(s, Inches(1.2), Inches(6.5), Inches(10), Inches(0.5),
             "MA Ecopsychology \u00b7 MS Creative Technology & Design (CU ATLAS) \u00b7 Founding engineer \u00d72 \u00b7 Ex-Google \u00b7 10+ years full stack",
             font_size=12, color=FG_DIM)


# ═══════════════════════════════════════════════
# SLIDE 9 — The Ask
# ═══════════════════════════════════════════════
s = new_slide()

section_label(s, "The Ask")
slide_headline(s, "Raising $300K to hire\nthe core team.")

# Ask details in a card
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(1.0), Inches(3.3), Inches(6), Inches(2.2))
shape.fill.solid()
shape.fill.fore_color.rgb = WHITE
shape.line.color.rgb = BORDER
shape.line.width = Pt(1)

ask_items = [
    ("Instrument", "SAFE (YC standard)"),
    ("Use of funds", "Core team full-time through 2026"),
    ("Goal", "Production launch by June \u00b7 revenue immediately"),
]

y = Inches(3.55)
for label, val in ask_items:
    add_text_box(s, Inches(1.4), y, Inches(2.2), Inches(0.4),
                 label, font_size=13, color=FG_DIM)
    add_text_box(s, Inches(3.6), y, Inches(3.2), Inches(0.4),
                 val, font_size=16, color=FG, font_name=FONT_SERIF)
    y += Inches(0.55)

# Closing line
add_text_box(s, Inches(1.2), Inches(5.8), Inches(10), Inches(1.0),
             "Open source. Local-first.\nBalance and choice \u2014 from the foundation.",
             font_size=28, color=FG, font_name=FONT_SERIF, line_spacing=1.3)

add_text_box(s, Inches(1.2), Inches(6.7), Inches(10), Inches(0.4),
             "Aaron Gabriel Neyer  \u00b7  aaron@parachute.computer  \u00b7  Boulder, CO",
             font_size=12, color=FG_DIM)


# ── Save ──
out_path = "/home/sandbox/parachute-computer/nvc-pitch-deck.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
